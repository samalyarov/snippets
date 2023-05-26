'''
A list of functions for automating reports with python-pptx
'''

# Importing libraries
import pandas as pd
import numpy as np

from pptx.chart.data import CategoryChartData, ChartData, BubbleChartData
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

# Functions for displaying graphs:

def plot_bubble(chart_position:str, 
                categories_names:list, 
                x_values:list,
                y_values:list,
                size_values:list,
                slides_data:pd.DataFrame,
                latest_date:str):
    
    '''
    A function for replacing bubble graphs data

    Args:
    chart_position (str): position of graph object on the slide (example - slide.shapes[0].chart)
    categories_names (list): a list with category names (how are the categories called - what to show on the graph)
    x_values (list): a list with values locations (in the slides_date df) for X axis
    y_values (list): a list with values locations (in the slides_date df) for Y axis
    size_values (list): a list with values locations (in the slides_date df) for bubble size
    slides_data(DataFrame): dataframe with data for the graph
    latest_date(str): latest report date (from which column to take data)
    '''
    
    # Selecting the chart
    chart = chart_position
    
    # Creating the chart data object
    chart_data = BubbleChartData()

    x_data = [round(slides_data.loc[x, latest_date]) for x in x_values]
    y_data = [round(slides_data.loc[y, latest_date]) for y in y_values]
    size_data = [round(slides_data.loc[z, latest_date]) for z in size_values]

    # Replacing data
    curr_series = chart_data.add_series('Slide data')
    for i, category in enumerate(categories_names):
        curr_series.add_data_point(x_data[i], y_data[i], size_data[i])
    chart.replace_data(chart_data)

    # Editing the data labels and font
    for i, category in enumerate(categories_names):
        chart.series[0].points[i].data_label.text_frame.text = category
        edit_text(text=chart.series[0].points[i].data_label.text_frame.paragraphs[0],
                  font_name='Calibri',
                  font_size=11,
                  font_italic=False,
                  font_bold=False,
                  font_shadow=False,
                  font_color_rgb=RGBColor(64,64,64))

        
    chart.replace_data(chart_data)



def plot_dynamic(chart_position:str, 
                 date_columns:list, 
                 num_lines:int, 
                 lines_locations:list, 
                 slides_data:pd.DataFrame,
                 lines_names:list,
                 abs_values=True,
                 rounding_level=0):
    '''
    A function for replacing most graphs data (no bubble|lines)

    Args:
    chart_position(str): position of graph object on the slide (example - slide.shapes[0].chart)
    date_columns(list): a list of date columns (from DF) from which to get data
    num_lines(int): number of lines in a graph
    lines_locations(list): a list with values locations (in the slides_date df), which rows to take
    slides_data(DataFrame): dataframe with data for the graph
    lines_names(list): a list with series names
    abs_values(bool): displaying absolute values as labels or not
    rounding_level(int): number of decimal places to display
    '''

    # Selecting the chart:
    chart = chart_position

    # Creating the chart data object:
    chart_data = ChartData()
    date_columns_full = []

    # Check which dates (columns) have data in them and only use those:
    for date in date_columns:
        if not pd.isna(slides_data.loc[lines_locations[0], date]):
            date_columns_full.append(date)
    chart_data.categories = date_columns_full

    # Creating a separate list for each line:
    lines_to_fill = [[] for i in range(num_lines)]

    # Adding tha data for each line into the corresponding list
    for date in date_columns_full:
        for i, line in enumerate(lines_to_fill):
            data_point = lines_locations[i]
            line.append(round(slides_data.loc[data_point, date], rounding_level))

    for i, name in enumerate(lines_names):
        chart_data.add_series(name, lines_to_fill[i])

    # Replacing data:
    chart.replace_data(chart_data)

    if abs_values:
        # Formatting labels if required:
        for chart_series in chart.series:
            chart_series.data_labels.number_format = '#;#'    

def plot_static(chart_position:int, 
                categories_names:list, 
                series_names:list, 
                series_data:list,
                slides_data:pd.DataFrame,
                latest_date:str,
                abs_values=True):
    '''
    A function for replacing regular chart data

    Args:
    chart_position (str): position of graph object on the slide (example - slide.shapes[0].chart)
    categories_names (list): a list with categories names
    series_names (list): a list with series names
    series_data(list of lists): a list of lists with df row numbers containing data for each category. One list - one category, with 1 entry for each series.
    slides_data(DataFrame): dataframe with data for the graph
    latest_date(str): last date (which column to use)
    abs_values(bool): displaying absolute values as labels or not
    '''

    
    # Selecting the cahrt:
    chart = chart_position

    # Gathering data:
    new_series_data = []
    for list in series_data:
        new_list = []
        for item in list:
            data_point = round(slides_data.loc[item, latest_date])
            new_list.append(data_point)
        new_series_data.append(new_list)


    # Creating the chart data object:
    chart_data = ChartData()
    chart_data.categories = categories_names

    # Collecting data for inserting:
    for i, u_series in enumerate(series_names):
        chart_data.add_series(u_series, new_series_data[i])

    # Replacing data:
    chart.replace_data(chart_data)

    if abs_values:
        # Formatting labels if required:
        for chart_series in chart.series:
            chart_series.data_labels.number_format = '#;#'
        
# Functions for editing text data:
        
def edit_text(text:object, 
              font_name:str, 
              font_size:int, 
              font_color_rgb:list,
              font_italic=False, 
              font_bold=False,
              font_shadow=False):
    '''
    Function for easier text editing

    Args:
    text (object): position of graph object on the slide (example - slide.shapes[0].text_frame.paragraphs[0])
    font_color_rgb(list): tuple or a list containing color RGB value (example - font_color_rgb=RGBColor(127,127,127))
    
    '''
    text.font.name = font_name
    text.font.size = Pt(font_size)
    text.font.italic = font_italic
    text.font.bold = font_bold
    text.font.color.rgb = font_color_rgb
    text.font.shadow = font_shadow

def calc_diff(position:int, 
              negative=False
              ) -> float:
    '''
    Function to calculate difference between the latest date and the one before that

    Args:
    position (int): which dataframe row corresponds to data of interest (where to calculate)
    negative (bool): whether the value is negative (if it is - simple substraction would not work)

    Returns:
    current_difference (float): difference between two values - from the last date and the one before that
    '''
    global slides_data  # Using global variable slides_data - the dataframe with all data
    global latest_date  
    global pre_latest_date 

    if negative == False:
        current_difference = round(slides_data.loc[position, latest_date]) - round(slides_data.loc[position, pre_latest_date])
    elif negative == True:
        current_difference = -1 * (round(slides_data.loc[position, latest_date]) - round(slides_data.loc[position, pre_latest_date]))
        
    return current_difference

def print_out_diff(current_diff:float, 
                   threshold:float, 
                   significant=False)-> str:
    """
    A function for creating a text containing the difference. Useful for easily creating text boxes indicating changes since last measurement. 

    Args:
    current_diff (float): difference between two values - from the last date and the one before that
    threshold (float): what difference level do we consider significant. 
    significant (bool): if the metric itself is considered significant. If it is - then any number is displayed, if not - only if the number is bigger then threshold

    Returns:
    curr_text (str): text itself with a sign (+ or -), value and 'p.p.' standing for 'percentage points'
    """
    if significant == False:
        if current_diff >= threshold:
            curr_text = f'+{current_diff} p.p.'
        elif current_diff <= -threshold:
            curr_text = f'{current_diff} p.p.'
        else:
            curr_text = ''
            
    if significant == True:
        if current_diff > 0:
            curr_text = f'+{current_diff} p.p.'
        elif current_diff == 0:
            curr_text = '0 p.p.'
        elif current_diff < 0:
            curr_text = f'{current_diff} p.p.'
            
    return curr_text
